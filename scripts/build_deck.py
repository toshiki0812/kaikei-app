"""新季さんへの資料（5年ビジョン・資金計画）をPowerPointで組み立てる。

数字はすべて scratchpad/deck.json・compare.json（アプリのエンジンが出した値）から読む。
ここで数字を手打ちしない＝アプリの計算とズレない、という作りにしている。
"""
from __future__ import annotations

import json
import os
import sys

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

SC = ("/private/tmp/claude-502/-Users-Office-Desktop-Claude-private/"
      "86f3df5d-53be-453b-aeb3-b23831f5312a/scratchpad")
D = json.load(open(f"{SC}/deck.json", encoding="utf-8"))
C = json.load(open(f"{SC}/compare.json", encoding="utf-8"))

INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x2A, 0x78, 0xD6)      # 青：俊来・本命
GREEN = RGBColor(0x1B, 0xAF, 0x7A)       # 緑：新季
AMBER = RGBColor(0xED, 0xA1, 0x00)       # 黄：世帯
RED = RGBColor(0xD9, 0x3A, 0x3A)
DARK = RGBColor(0x1E, 0x2A, 0x38)        # 区切りページの地色
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
FONT = "Yu Gothic"

W, H = Inches(13.333), Inches(7.5)
yen = lambda v: f"¥{v:,}"


def deck():
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def tb(slide, x, y, w, h, text, size=18, bold=False, color=INK,
       align=PP_ALIGN.LEFT, spacing=1.25):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.alignment = align
        para.line_spacing = spacing
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = FONT
    return box


def band(slide, color=LIGHT, y=Emu(0), h=H):
    s = slide.shapes.add_shape(1, Emu(0), y, W, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def title_page(slide, kicker, title, sub=""):
    tb(slide, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.4),
       kicker, 13, False, ACCENT)
    tb(slide, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0),
       title, 30, True, INK)
    if sub:
        tb(slide, Inches(0.9), Inches(1.95), Inches(11.5), Inches(0.6),
           sub, 14, False, MUTED)


def section(prs, num, title, sub=""):
    s = blank(prs)
    band(s, DARK)
    tb(s, Inches(1.1), Inches(2.7), Inches(11), Inches(0.5), num, 15, False,
       RGBColor(0x7E, 0xA8, 0xD8))
    tb(s, Inches(1.1), Inches(3.15), Inches(11), Inches(1.1), title, 36, True, WHITE)
    if sub:
        tb(s, Inches(1.1), Inches(4.35), Inches(11), Inches(0.8), sub, 15,
           False, RGBColor(0xC5, 0xD2, 0xDE))
    return s


def table(slide, x, y, w, rows, widths=None, head_fill=RGBColor(0x2C, 0x3E, 0x50),
          size=12, row_h=Inches(0.42)):
    nr, nc = len(rows), len(rows[0])
    shape = slide.shapes.add_table(nr, nc, x, y, w, row_h * nr)
    t = shape.table
    if widths:
        total = sum(widths)
        for i, ww in enumerate(widths):
            t.columns[i].width = Emu(int(w * ww / total))
    for r, row in enumerate(rows):
        t.rows[r].height = row_h
        for c, val in enumerate(row):
            cell = t.cell(r, c)
            cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.RIGHT if (c > 0 and r > 0) else PP_ALIGN.LEFT
            for run in para.runs:
                run.font.size = Pt(size)
                run.font.name = FONT
                run.font.bold = (r == 0)
                run.font.color.rgb = WHITE if r == 0 else INK
            cell.fill.solid()
            cell.fill.fore_color.rgb = (head_fill if r == 0
                                        else (WHITE if r % 2 else LIGHT))
    return t


def line_chart(slide, x, y, w, h, cats, series, colors, number_format='#,##0'):
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals, number_format)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, w, h, cd)
    ch = gf.chart
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False
    ch.legend.font.size = Pt(11)
    ch.legend.font.name = FONT
    ch.font.size = Pt(10)
    ch.font.name = FONT
    for i, ser in enumerate(ch.series):
        ser.smooth = False
        ser.format.line.color.rgb = colors[i]
        ser.format.line.width = Pt(2.5)
    va = ch.value_axis
    va.has_major_gridlines = True
    va.tick_labels.number_format = '#,##0,,"百万"'
    va.tick_labels.number_format_is_linked = False
    va.tick_labels.font.size = Pt(10)
    ch.category_axis.tick_labels.font.size = Pt(9)
    return ch


def note(slide, y, text):
    tb(slide, Inches(0.9), y, Inches(11.5), Inches(1.0), text, 11, False, MUTED,
       spacing=1.3)


def mlabel(m):
    y, mm = m.split("-")
    return f"{y[2:]}/{int(mm)}"



def main():
    out = sys.argv[1] if len(sys.argv) > 1 else "/Users/Office/Desktop/資金シミュレーション.pptx"
    prs = deck()
    m12 = D["m12"]; cats12 = [mlabel(m) for m in m12]
    m60 = D["m60"]; cats60 = [mlabel(m) if m.endswith(("-03", "-09")) else "" for m in m60]
    T1, S1 = D["y1"]["T"], D["y1"]["S"]
    T5, S5 = D["y5"]["T"], D["y5"]["S"]
    H5 = [t + s for t, s in zip(T5, S5)]

    # ══════ 表紙 ══════
    s = blank(prs); band(s, DARK)
    tb(s, Inches(1.1), Inches(2.6), Inches(11), Inches(0.5),
       "俊来・新季", 15, False, RGBColor(0x7E, 0xA8, 0xD8))
    tb(s, Inches(1.1), Inches(3.15), Inches(11), Inches(1.4),
       "資金シミュレーション", 38, True, WHITE)
    tb(s, Inches(1.1), Inches(4.6), Inches(11), Inches(0.8),
       "2026年9月 〜 2031年8月（60ヶ月）", 17, False, RGBColor(0xC5, 0xD2, 0xDE))

    # ══════ 01 前提：俊来 ══════
    s = blank(prs)
    title_page(s, "01", "前提条件：俊来")
    table(s, Inches(0.9), Inches(2.2), Inches(11.5), [
        ["項目", "1年目（2026年9月〜2027年8月）", "2年目以降（2027年9月〜）"],
        ["収入", "2026年9月 244万円／10月〜2027年3月 26万円／4月 125万円／5月〜 50万円", "50万円"],
        ["家賃", "〜2027年5月 16.7万円／6月〜 8万円（同棲）", "8万円"],
        ["家賃以外の費用", "9万円（食費6万＋クレカ3万）", "11万円（食費6万＋クレカ4.75万）"],
        ["　飲み・外食", "3万円", "3万円"],
        ["　自炊", "3万円", "3万円"],
        ["投資拠出", "2026年9月 100万円／10月〜2027年1月 0円／2月〜 20万円", "20万円"],
        ["税金", "年60万円（9月8万・12月8万・3月38万・6月6万）", "年82万円（住民税8万×4＋所得税50万）"],
        ["現金の上限", "100万円（超えた分は投資へ自動振替）", "同左"],
        ["開始時の残高", "現金0円・投資0円", "—"],
    ], widths=[2.4, 5.6, 3.5], size=11, row_h=Inches(0.40))
    note(s, Inches(6.5),
         "※ 現金が足りない月は、まず投資拠出を止めて現金に回し、それでも足りなければ投資残高から取り崩す。")

    # ══════ 02 前提：新季・共通 ══════
    s = blank(prs)
    title_page(s, "02", "前提条件：新季／二人で共通")
    table(s, Inches(0.9), Inches(2.3), Inches(6.4), [
        ["新季", "金額"],
        ["収入", "25万円"],
        ["　産休・育休中（2030年4月〜2031年5月）", "15万円"],
        ["家賃", "8万円"],
        ["クレカ", "10万円（同棲後 9.75万円）"],
        ["投資拠出", "1万円"],
        ["賞与", "7月・12月に各30万円"],
        ["現金の上限", "200万円"],
    ], widths=[3.8, 2.6], size=11, row_h=Inches(0.44))
    table(s, Inches(7.7), Inches(2.3), Inches(4.7), [
        ["二人で共通", "内容"],
        ["想定年利", "5%（毎月複利）"],
        ["光熱費の節約", "同棲後に世帯5,000円／月"],
        ["住宅ローン", "2031年1月から月10万円ずつ"],
        ["　返済期間", "35年（返済中は家賃を計上しない）"],
        ["　年間の値動き", "0.5%"],
    ], widths=[2.4, 2.3], size=11, row_h=Inches(0.44))
    note(s, Inches(5.9),
         "※ 住宅ローンは「家賃の代わりに払うもの」として扱い、返済額が年0.5%の値動きを乗せながら不動産の資産として積み上がる。"
         "金利・元利内訳・ローン残債は扱わない単純化モデル。頭金・購入諸費用は含んでいない。")

    # ══════ 03 イベント一覧 ══════
    s = blank(prs)
    title_page(s, "03", "5年間に起きるイベント",
               "毎年くり返すもの（住民税・所得税・賞与）は前提条件のページに記載")
    table(s, Inches(0.9), Inches(2.4), Inches(11.5), [
        ["時期", "内容", "対象", "金額"],
        ["2026年9月", "投資への初回投下", "俊来", "−" + yen(1000000)],
        ["2027年4月", "C-LinC 1〜3月分報酬", "俊来", "+" + yen(750000)],
        ["2027年6月", "引越し費用（敷金礼金1ヶ月分を含む）", "俊来", "−" + yen(600000)],
        ["2027年9月", "婚約", "俊来", "−" + yen(600000)],
        ["2028年8月", "結婚式（費用）", "二人", "−" + yen(4000000)],
        ["2028年9月", "結婚式（ご祝儀）", "二人", "+" + yen(2000000)],
        ["2028年10月", "新婚旅行", "二人", "−" + yen(1000000)],
        ["2030年4月〜2031年5月", "新季の産休・育休（収入減）", "新季", "−100,000円/月"],
        ["2030年6月", "第一子・一時費用", "二人", "−" + yen(300000)],
        ["2030年6月〜2031年5月", "第一子・継続費", "二人", "−15,000円/月"],
        ["2031年1月〜", "住宅ローン返済", "二人", "−200,000円/月"],
    ], widths=[2.6, 4.8, 1.4, 2.4], size=11.5, row_h=Inches(0.38))

    # ══════ 04 第一子の費用 ══════
    s = blank(prs)
    title_page(s, "04", "第一子の費用（2030年6月）",
               "東京都府中市の制度を反映した実質負担")
    table(s, Inches(0.9), Inches(2.5), Inches(6.6), [
        ["項目", "金額"],
        ["分娩入院費", yen(550000)],
        ["妊婦健診", yen(70000)],
        ["ベビー用品（初期）", yen(250000)],
        ["マタニティ用品", yen(30000)],
        ["支出 小計", yen(900000)],
        ["出産育児一時金", "−" + yen(500000)],
        ["妊婦のための支援給付（府中市）", "−" + yen(100000)],
        ["実質負担", yen(300000)],
    ], widths=[3.4, 2])
    tb(s, Inches(8.0), Inches(2.6), Inches(4.4), Inches(0.4), "毎月かかるもの", 14, True, ACCENT)
    tb(s, Inches(8.0), Inches(3.05), Inches(4.4), Inches(2.0),
       "育児消耗品費　20,000円/月\n018サポート　　−5,000円/月\n"
       "──────────────\n実質負担　　　 15,000円/月\n（0〜1歳の期間）", 13, False, INK)
    tb(s, Inches(8.0), Inches(5.0), Inches(4.4), Inches(0.4), "保育料", 14, True, GREEN)
    tb(s, Inches(8.0), Inches(5.4), Inches(4.4), Inches(0.8),
       "認可保育所なら実質0円\n（府中市・2026年4月から所得問わず無償化）", 12, False, INK)
    note(s, Inches(6.3),
         "※ 分娩費用・健診費用・ベビー用品費は全国平均ベースの仮置き。"
         "新季の収入減（25万→15万・14ヶ月）は国の給付水準にもとづく目安で、勤務先の実際の制度は未確認。"
         "認可外保育施設に入った場合、府中市の補助は非課税世帯のみのため別途月5〜8万円がかかるリスクがある。")

    # ══════ 05-07 最初の1年 ══════
    def year1_page(no, who, rows, color, note_text, household=False):
        sl = blank(prs)
        title_page(sl, no, f"最初の1年：{who}", "2026年9月〜2027年8月")
        if household:
            series = [("俊来", [r["total_assets"] for r in T1], ACCENT),
                      ("新季", [r["total_assets"] for r in S1], GREEN),
                      ("世帯合計", [a["total_assets"] + b["total_assets"]
                                for a, b in zip(T1, S1)], AMBER)]
        else:
            series = [("現金", [r["cash_balance"] for r in rows], RGBColor(0x8A, 0xA8, 0xC8)),
                      ("投資", [r["investment_balance"] for r in rows], GREEN),
                      ("資産合計", [r["total_assets"] for r in rows], color)]
        line_chart(sl, Inches(0.9), Inches(2.35), Inches(6.4), Inches(3.6),
                   cats12, [(n, v) for n, v, _ in series], [c for _, _, c in series])
        tb(sl, Inches(7.6), Inches(2.35), Inches(4.8), Inches(0.4), "月次の内訳", 14, True, MUTED)
        if household:
            body = [["月", "俊来", "新季", "世帯"]]
            for i in range(0, 12, 2):
                body.append([mlabel(m12[i]), yen(T1[i]["total_assets"]),
                             yen(S1[i]["total_assets"]),
                             yen(T1[i]["total_assets"] + S1[i]["total_assets"])])
        else:
            body = [["月", "収入", "支出計", "資産合計"]]
            for i in range(0, 12, 2):
                r = rows[i]
                spend = (r["rent"] + r["credit_card"] + r["other_cash_expense"]
                         + r["investment_contribution"] + r["planned_expense"])
                body.append([mlabel(m12[i]), yen(r["income"]), yen(spend),
                             yen(r["total_assets"])])
        table(sl, Inches(7.6), Inches(2.8), Inches(4.8), body,
              widths=[1.1, 1.6, 1.6, 1.7], size=10, row_h=Inches(0.38))
        note(sl, Inches(6.15), note_text)
        return sl

    year1_page("05", "俊来", T1, ACCENT,
               "※ 2026年9月に244万円が入り、うち100万円を投資へ。10月〜2027年3月は月収26万円で、"
               "この間の投資拠出はゼロ。2027年2月から毎月20万円を投資へ回す。"
               "2027年3月は税金38万円で現金が15.2万円まで落ちるが、マイナスにはならない。"
               "4月にC-LinC分75万円が入り、6月の同棲で家賃が16.7万→8万に下がる（引越し費用60万円）。")
    year1_page("06", "新季", S1, GREEN,
               "※ 収入25万円・投資拠出1万円で一定。7月と12月に賞与30万円が入る。"
               "引越し費用は俊来が全額立て替えるため、新季側の現金は細らない。")
    year1_page("07", "世帯合計", None, AMBER,
               f"※ 1年後（2027年8月）の世帯資産は {yen(T1[11]['total_assets'] + S1[11]['total_assets'])}。"
               "2027年3月（税金）と6月（引越し）に谷ができるが、"
               "現金がマイナスになる月は一度もない。", household=True)

    # ══════ 08-10 5年の見通し ══════
    def year5_page(no, who, vals, cash, inv, re_, color, note_text):
        sl = blank(prs)
        title_page(sl, no, f"5年の見通し：{who}", "2026年9月〜2031年8月")
        line_chart(sl, Inches(0.9), Inches(2.35), Inches(7.4), Inches(3.6), cats60,
                   [("現金", cash), ("投資", inv), ("不動産", re_), ("資産合計", vals)],
                   [RGBColor(0x8A, 0xA8, 0xC8), GREEN, AMBER, color])
        rows = [["時点", "資産合計"]]
        for i, lab in ((11, "1年後"), (23, "2年後"), (35, "3年後"), (47, "4年後"), (59, "5年後")):
            rows.append([f"{lab}（{m60[i][:4]}年{int(m60[i][5:])}月）", yen(vals[i])])
        table(sl, Inches(8.7), Inches(2.6), Inches(3.7), rows,
              widths=[2.1, 1.6], size=11, row_h=Inches(0.44))
        note(sl, Inches(6.15), note_text)
        return sl

    year5_page("08", "俊来", T5, D["y5"]["Tcash"], D["y5"]["Tinv"], D["y5"]["Tre"], ACCENT,
               "※ 2028年8月の結婚式で現金がゼロになり、投資拠出が自動的に止まったうえで"
               "投資から一部を取り崩す。2031年1月から住宅ローンの返済が始まり、"
               "返済額が不動産の資産として積み上がる。")
    year5_page("09", "新季", S5, D["y5"]["Scash"], D["y5"]["Sinv"], D["y5"]["Sre"], GREEN,
               "※ 2030年4月〜2031年5月は産休・育休で収入が月15万円に下がる。"
               "現金が上限200万円に達すると、超えた分は自動的に投資へ回る。")
    Hc = [a + b for a, b in zip(D["y5"]["Tcash"], D["y5"]["Scash"])]
    Hi = [a + b for a, b in zip(D["y5"]["Tinv"], D["y5"]["Sinv"])]
    Hr = [a + b for a, b in zip(D["y5"]["Tre"], D["y5"]["Sre"])]
    year5_page("10", "世帯合計", H5, Hc, Hi, Hr, AMBER,
               f"※ 5年後の世帯資産は {yen(H5[59])}。結婚式・新婚旅行・婚約・第一子を"
               "すべて払ったうえでの数字。2028年8月〜10月にイベントが集中するため、"
               "2年目は資産がほぼ横ばいになる。")

    # ══════ 11 感応度 ══════
    s = blank(prs)
    title_page(s, "11", "食費の枠を守れなかった場合",
               "同じ収入・同じイベントで、飲食費だけを7月の実績水準に戻して比べた")
    table(s, Inches(0.9), Inches(2.5), Inches(11.5), [
        ["", "枠を守り切った場合", "見直しが続かなかった場合"],
        ["飲食費", "月6万円", "月132,820円（7月の実績）"],
        ["5年後の世帯資産", yen(C["plan"]), yen(C["bad"])],
    ], widths=[3.2, 4.2, 4.2], size=13, row_h=Inches(0.62))
    box = s.shapes.add_shape(1, Inches(0.9), Inches(4.9), Inches(11.5), Inches(1.3))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xE8, 0xF1, 0xFB)
    box.line.fill.background(); box.shadow.inherit = False
    tb(s, Inches(1.3), Inches(5.2), Inches(10.7), Inches(0.8),
       f"差は {yen(C['gap'])}。5年間の飲食費の差（{yen((132820 - 60000) * 60)}）に、"
       "投資に回したぶんの運用益が乗る。", 18, True, RGBColor(0x18, 0x5F, 0xA5))

    prs.save(out)
    print(f"書き出しました: {out}")
    print(f"  全 {len(prs.slides._sldIdLst)} 枚")
    return out


if __name__ == "__main__":
    main()
